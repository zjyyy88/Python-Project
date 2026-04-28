from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TILE_BG = RGBColor(219, 130, 226)
BORDER = RGBColor(20, 20, 20)
TEXT = RGBColor(0, 0, 0)


ELEMENTS = [
    {
        "symbol": "F",
        "name": "Fluorine",
        "mass": "18.998403",
        "z": "9",
        "ie": "1681.0",
        "en": "3.98",
        "ox": ["-1"],
        "config": "[He] 2s2 2p5",
    },
    {
        "symbol": "Cl",
        "name": "Chlorine",
        "mass": "35.453",
        "z": "17",
        "ie": "1251.2",
        "en": "3.16",
        "ox": ["+7", "+5", "+3", "+1", "-1"],
        "config": "[Ne] 3s2 3p5",
    },
    {
        "symbol": "Br",
        "name": "Bromine",
        "mass": "79.904",
        "z": "35",
        "ie": "1139.9",
        "en": "2.96",
        "ox": ["+7", "+5", "+3", "+1", "-1"],
        "config": "[Ar] 3d10 4s2 4p5",
    },
    {
        "symbol": "I",
        "name": "Iodine",
        "mass": "126.9044",
        "z": "53",
        "ie": "1008.4",
        "en": "2.66",
        "ox": ["+7", "+5", "+1", "-1"],
        "config": "[Kr] 4d10 5s2 5p5",
    },
]


def put_text(slide, x, y, w, h, text, size, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = TEXT


def draw_tile(slide, x, y, w, h, data):
    tile = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    tile.fill.solid()
    tile.fill.fore_color.rgb = TILE_BG
    tile.line.color.rgb = BORDER
    tile.line.width = Pt(1.5)

    put_text(slide, x + 0.10, y + 0.05, 1.45, 0.24, data["mass"], 15)
    put_text(slide, x + 2.23, y + 0.02, 0.65, 0.32, data["z"], 24, bold=False, align=PP_ALIGN.RIGHT)
    put_text(slide, x + 0.10, y + 0.31, 1.00, 0.20, data["ie"], 11)
    put_text(slide, x + 1.48, y + 0.31, 0.78, 0.20, data["en"], 11, align=PP_ALIGN.RIGHT)

    ox = "\n".join(data["ox"])
    put_text(slide, x + 2.37, y + 0.56, 0.48, 0.95, ox, 9, align=PP_ALIGN.RIGHT)

    symbol_size = 56 if len(data["symbol"]) == 1 else 46
    put_text(slide, x + 0.10, y + 0.54, 1.45, 0.90, data["symbol"], symbol_size)
    put_text(slide, x + 0.10, y + 1.88, 1.95, 0.30, data["name"], 17)
    put_text(slide, x + 0.10, y + 2.16, 2.55, 0.30, data["config"], 13)


def main():
    prs = Presentation()
    prs.slide_width = Inches(3.0)
    prs.slide_height = Inches(12.4)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tile_h = 3.1
    for idx, data in enumerate(ELEMENTS):
        draw_tile(slide, x=0.0, y=idx * tile_h, w=3.0, h=tile_h, data=data)

    out_path = Path(__file__).resolve().parent / "halogen_column_vector.pptx"
    prs.save(out_path)
    print(f"Generated: {out_path}")


if __name__ == "__main__":
    main()